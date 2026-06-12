"""PPTX to Markdown conversion using python-pptx."""

import html
import io
import logging

from pptx import Presentation

from app.config import settings
from app.converters.base import clean_lines
from app.models.response import SlideJson

logger = logging.getLogger(__name__)


def _extract_text_from_table(table) -> str:
    """Extract text from a PowerPoint table and convert to Markdown table.

    Args:
        table: A python-pptx table object.

    Returns:
        A Markdown string representing the table.
    """
    if not table.rows or not table.columns:
        return ""

    lines = []
    rows = []
    for row in table.rows:
        row_values = []
        for cell in row.cells:
            text = html.escape(cell.text).replace("|", "\\|").replace("\n", " ")
            row_values.append(text)
        rows.append(row_values)

    if not rows:
        return ""

    # Header row
    header = "| " + " | ".join(rows[0]) + " |"
    lines.append(header)

    # Separator row
    separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
    lines.append(separator)

    # Data rows
    for row in rows[1:]:
        line = "| " + " | ".join(row) + " |"
        lines.append(line)

    return "\n".join(lines)


def _extract_slide_text(slide, slide_number: int) -> list[str]:
    """Extract text content from a single slide.

    Args:
        slide: A python-pptx slide object.
        slide_number: The 1-based slide number.

    Returns:
        A list of Markdown lines representing the slide content.
    """
    lines = []

    # Handle slide title (escaped to prevent XSS)
    if slide.shapes.title:
        title_text = html.escape(slide.shapes.title.text).strip()
        if title_text:
            lines.append(f"## {title_text}")
        else:
            lines.append(f"## Slide {slide_number}")
    else:
        lines.append(f"## Slide {slide_number}")

    # Process all shapes
    for shape in slide.shapes:
        # Skip title (already handled)
        if shape == slide.shapes.title:
            continue

        # Handle tables
        if shape.has_table:
            table_md = _extract_text_from_table(shape.table)
            if table_md:
                lines.append("")
                lines.append(table_md)
                continue

        # Handle text frames
        if shape.has_text_frame:
            text_frame = shape.text_frame

            for paragraph in text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue

                # Detect bullet points (escaped to prevent XSS)
                escaped_text = html.escape(text)
                if paragraph.level > 0:
                    indent = "  " * paragraph.level
                    lines.append(f"{indent}- {escaped_text}")
                else:
                    # Check if the paragraph is a bullet
                    try:
                        if paragraph.style is not None:
                            # python-pptx doesn't have a direct bullet property
                            # but we can check the paragraph's level
                            lines.append(f"- {escaped_text}" if paragraph.level >= 0 else escaped_text)
                        else:
                            lines.append(escaped_text)
                    except Exception:
                        lines.append(escaped_text)

    return lines


def convert_pptx_to_md(file_bytes: bytes) -> str:
    """Pure Python PPTX → MD conversion (python-pptx).

    Converts all slides to Markdown, handling:
    - Slide titles
    - Text frames with paragraphs
    - Bullet points
    - Tables within slides
    - Empty slides

    Args:
        file_bytes: The raw bytes of the PPTX file.

    Returns:
        A Markdown string representing the entire presentation.

    Raises:
        Exception: If the file is not a valid PPTX.
    """
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        logger.error("Invalid PPTX file: %s", e)
        raise

    # Log presentation properties
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    logger.debug(
        "PPTX properties: slides=%d, width=%d, height=%d",
        len(prs.slides),
        slide_width,
        slide_height,
    )

    all_lines = []
    slide_count = len(prs.slides)

    for idx, slide in enumerate(prs.slides, start=1):
        try:
            slide_lines = _extract_slide_text(slide, idx)
            all_lines.extend(slide_lines)

            # Add separator between slides (but not after the last one)
            if idx < slide_count:
                all_lines.append("")
                all_lines.append("---")
                all_lines.append("")
        except Exception as e:
            logger.warning("Error converting slide %d: %s", idx, e)
            all_lines.append(f"## Slide {idx}")
            all_lines.append(f"_(error converting slide: {e})_")

            if idx < slide_count:
                all_lines.append("")
                all_lines.append("---")
                all_lines.append("")

    # Check for unsupported elements and log warnings
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            # Check for charts
            if shape.has_chart:
                logger.warning(
                    "Slide %d contains a chart which is not supported in Markdown conversion",
                    idx,
                )
            # Check for media (images, videos, audio)
            if shape.shape_type in (13, 14, 15):  # Picture, Media, Video
                logger.warning(
                    "Slide %d contains media content which is not supported in Markdown conversion",
                    idx,
                )

    # Post-process with clean_lines from base module
    result_lines = clean_lines(all_lines)

    return "\n".join(result_lines)


def _extract_pptx_content(file_bytes: bytes) -> list[tuple[int, str, list[str]]]:
    """Extract PPTX content as slides.

    Args:
        file_bytes: Raw PPTX file bytes.

    Returns:
        List of tuples (slide_num, title, text_lines).
    """
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        logger.error("Invalid PPTX file: %s", e)
        raise

    results = []

    for idx, slide in enumerate(prs.slides, start=1):
        try:
            lines = _extract_slide_text(slide, idx)
            # Get title or default to slide number
            if slide.shapes.title:
                title_text = html.escape(slide.shapes.title.text).strip()
                title = title_text if title_text else f"Slide {idx}"
            else:
                title = f"Slide {idx}"

            results.append((idx, title, lines))
        except Exception as e:
            logger.warning("Error extracting slide %d: %s", idx, e)
            results.append((idx, f"Slide {idx}", [f"_(error converting slide: {e})_"]))

    return results


def convert_pptx_to_json(file_bytes: bytes) -> dict:
    """Convert PPTX to JSON format.

    Args:
        file_bytes: Raw PPTX file bytes.

    Returns:
        Dict with slides and metadata in JSON structure.
    """
    results = _extract_pptx_content(file_bytes)

    return {
        "format": "pptx",
        "slides": [
            SlideJson(
                index=slide[0],
                title=None,
                content=[line for line in slide[2] if line.strip()]
            ).model_dump()
            for slide in results
        ],
        "metadata": {
            "format": "pptx",
            "size_bytes": len(file_bytes),
            "slide_count": len(results),
        },
    }


def convert_pptx_to_jsonl(file_bytes: bytes) -> str:
    """Convert PPTX to JSONL format with chunking.

    Args:
        file_bytes: Raw PPTX file bytes.

    Returns:
        JSONL string with start, chunk, and end events.
    """
    results = _extract_pptx_content(file_bytes)

    return _to_jsonl(results)


def _to_jsonl(results: list[tuple[int, str, list[str]]]) -> str:
    """Convert extraction results to JSONL format with chunking.

    Args:
        results: List of (slide_num, title, text_lines) tuples.

    Returns:
        JSONL string with start, chunk, and end events.
    """
    import json

    lines = []

    # Start event
    lines.append(json.dumps({
        "type": "start",
        "format": "pptx",
    }))

    # Convert to text representation for chunking
    all_text = []
    for slide_num, title, slide_lines in results:
        all_text.append(f"Slide {slide_num}: {title}")
        all_text.extend(slide_lines)

    chunk_size = settings.jsonl_chunk_size

    if all_text:
        chunks = [all_text[i:i + chunk_size] for i in range(0, len(all_text), chunk_size)]

        for chunk in chunks:
            lines.append(json.dumps({
                "type": "chunk",
                "content": "\n".join(chunk),
            }))

    # End event
    lines.append(json.dumps({
        "type": "end",
        "format": "pptx",
        "total_slides": len(results),
    }))

    return "\n".join(lines)
