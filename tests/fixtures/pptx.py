"""Fixtures for PPTX tests."""

__all__ = [
    "sample_pptx_bytes",
    "sample_pptx_with_table_bytes",
    "sample_pptx_empty_slide_bytes",
]

import io
import pytest


@pytest.fixture
def sample_pptx_bytes() -> bytes:
    """Generate a minimal PPTX presentation in memory using python-pptx."""
    from pptx import Presentation

    prs = Presentation()

    # Slide 1: title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Présentation de Test"
    slide.placeholders[1].text = "Sous-titre de la présentation"

    # Slide 2: title + content with bullets
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Deuxième Slide"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Premier point"
    p = tf.add_paragraph()
    p.text = "Deuxième point"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Sous-point"
    p.level = 1

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_with_table_bytes() -> bytes:
    """Generate a PPTX presentation containing a table in memory using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # Slide with table
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_box.text = "Slide avec Tableau"

    # Add table
    rows = 3
    cols = 3
    left = Inches(1)
    top = Inches(1)
    width = Inches(6)
    height = Inches(2)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.cell(0, 0).text = "En-tête 1"
    table.cell(0, 1).text = "En-tête 2"
    table.cell(0, 2).text = "En-tête 3"
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "B"
    table.cell(1, 2).text = "C"
    table.cell(2, 0).text = "1"
    table.cell(2, 1).text = "2"
    table.cell(2, 2).text = "3"

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_empty_slide_bytes() -> bytes:
    """Generate a PPTX presentation containing an empty slide in memory using python-pptx."""
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    prs.slides.add_slide(slide_layout)
    # No content added — empty slide

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
