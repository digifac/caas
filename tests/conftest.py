"""Shared fixtures for all tests."""

import asyncio
import io
import os
from typing import Any, Dict, List

import httpx
import pytest
from app.api import app
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas


@pytest.fixture
def clean_caas_env(monkeypatch: Any) -> None:
    """Remove all CAAS_ environment variables for testing defaults."""
    keys_to_remove = [k for k in os.environ if k.startswith("CAAS_")]
    for key in keys_to_remove:
        monkeypatch.delenv(key, raising=False)


@pytest.fixture(autouse=True)
def disable_rate_limiting():
    """Disable rate limiting during all tests."""
    app.state.rate_limiter.enabled = False
    yield


@pytest.fixture(autouse=True)
def reset_task_manager():
    """Reset the TaskManager between each test to avoid task accumulation.
    Also restores the original task_manager if a test replaced app.state.task_manager.
    """
    from app.api import app as main_app

    original_tm = main_app.state.task_manager
    original_tm._tasks.clear()
    original_tm._async_tasks.clear()
    original_tm._batches.clear()
    # Reset the semaphore so it's available again
    original_tm._semaphore = asyncio.Semaphore(original_tm._max_concurrent)
    yield
    # Restore original task_manager in case a test replaced it
    main_app.state.task_manager = original_tm
    original_tm._tasks.clear()
    original_tm._async_tasks.clear()
    original_tm._batches.clear()
    original_tm._semaphore = asyncio.Semaphore(original_tm._max_concurrent)


@pytest.fixture
async def async_client():
    """Client HTTP asynchrone pour tester l'API FastAPI."""
    async with httpx.AsyncClient(
        transport=httpx.ASGITransport(app=app), base_url="http://test"
    ) as client:
        yield client


@pytest.fixture
def sample_pdf_bytes() -> bytes:
    """Generate a minimal PDF with text in memory (zero-disk)."""
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    c.setFont("Helvetica", 12)
    c.drawString(72, 750, "Hello World")
    c.drawString(72, 730, "Ligne deux")
    c.save()
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_pdf_with_link_bytes() -> bytes:
    """Generate a PDF with a hyperlink in memory."""
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    c.setFont("Helvetica", 12)
    c.drawString(72, 750, "Visitez notre site")
    # Add an URI link with reportlab
    c.linkURL("https://example.com", (72, 740, 200, 760), thickness=0)  # type: ignore[attr-defined]
    c.save()
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_scanned_pdf_bytes() -> bytes:
    """Generate a 'scanned' PDF (image with French text, no text layer)."""
    import io as _io
    import os

    from PIL import Image, ImageDraw, ImageFont

    # Create a larger image with a TrueType font for reliable OCR
    img = Image.new("RGB", (1200, 300), color="white")
    draw = ImageDraw.Draw(img)

    # Use the arial.ttf font provided in the tests folder
    font_path = os.path.join(os.path.dirname(__file__), "arial.ttf")
    font = ImageFont.truetype(font_path, 36)

    # Draw black text on white background (simulates a scanned document)
    draw.text((40, 100), "Document scanné en français", fill="black", font=font)
    img_bytes = _io.BytesIO()
    img.save(img_bytes, format="PNG")
    img_bytes.seek(0)

    # Create a PDF with the image using reportlab (no text layer)
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    from reportlab.lib.utils import ImageReader

    c.drawImage(ImageReader(img_bytes), 50, 50, width=495, height=200)  # type: ignore[attr-defined]
    c.save()
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_docx_bytes() -> bytes:
    """
    Generate a minimal DOCX in memory.
    Uses the raw XML structure of a .docx (zip) without external dependencies.
    """
    import zipfile

    content_types = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>'
    relationships = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>'
    document_xml = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body><w:p><w:r><w:t>Bonjour le monde</w:t></w:r></w:p><w:p><w:r><w:t>Deuxième paragraphe</w:t></w:r></w:p></w:body></w:document>'

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", relationships)
        zf.writestr("word/document.xml", document_xml)

    return buf.getvalue()


@pytest.fixture
def sample_html_bytes() -> bytes:
    """
    Generate a minimal HTML document in memory with various elements.
    Includes headings, paragraphs, links, lists, bold/italic, and a table.
    """
    html = """<!DOCTYPE html>
<html>
<head><title>Test Page</title></head>
<body>
    <h1>Main Title</h1>
    <p>This is a <strong>bold</strong> and <em>italic</em> paragraph.</p>
    <p>Visit <a href="https://example.com">our website</a> for more info.</p>
    <ul>
        <li>First item</li>
        <li>Second item</li>
    </ul>
    <ol>
        <li>Step one</li>
        <li>Step two</li>
    </ol>
    <table>
        <tr><th>Name</th><th>Value</th></tr>
        <tr><td>A</td><td>1</td></tr>
        <tr><td>B</td><td>2</td></tr>
    </table>
    <blockquote>
        <p>A quoted paragraph</p>
    </blockquote>
    <pre><code>def hello():\n    print("world")</code></pre>
    <hr>
    <p>Final paragraph with an <img src="image.png" alt="test image"> inline.</p>
</body>
</html>"""
    return html.encode("utf-8")


@pytest.fixture
def sample_html_minimal_bytes() -> bytes:
    """Generate a minimal HTML document with just a heading and paragraph."""
    html = """<!DOCTYPE html>
<html>
<body>
    <h1>Hello World</h1>
    <p>This is a simple paragraph.</p>
</body>
</html>"""
    return html.encode("utf-8")


@pytest.fixture
def sample_html_latin1_bytes() -> bytes:
    """Generate an HTML document encoded in Latin-1."""
    html = """<!DOCTYPE html>
<html>
<body>
    <h1>Document en français</h1>
    <p>Des caractères spéciaux : àéîôù</p>
</body>
</html>"""
    return html.encode("latin-1")


@pytest.fixture
def sample_odt_bytes() -> bytes:
    """Generate a minimal ODT document in memory using raw XML (zip-based)."""
    import zipfile

    mimetype = "application/vnd.oasis.opendocument.text"
    manifest = """<?xml version="1.0" encoding="UTF-8"?>
<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0" manifest:version="1.2">
  <manifest:file-entry manifest:full-path="/" manifest:media-type="application/vnd.oasis.opendocument.text"/>
  <manifest:file-entry manifest:full-path="content.xml" manifest:media-type="text/xml"/>
</manifest:manifest>"""
    content_xml = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" office:version="1.2">
  <office:body>
    <office:text>
      <text:p>Titre du document</text:p>
      <text:p>Premier paragraphe</text:p>
      <text:p>Deuxième paragraphe</text:p>
    </office:text>
  </office:body>
</office:document-content>"""

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", mimetype)
        zf.writestr("META-INF/manifest.xml", manifest)
        zf.writestr("content.xml", content_xml)
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_odt_with_list_bytes() -> bytes:
    """Generate an ODT document with list items using raw XML (zip-based)."""
    import zipfile

    mimetype = "application/vnd.oasis.opendocument.text"
    manifest = """<?xml version="1.0" encoding="UTF-8"?>
<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0" manifest:version="1.2">
  <manifest:file-entry manifest:full-path="/" manifest:media-type="application/vnd.oasis.opendocument.text"/>
  <manifest:file-entry manifest:full-path="content.xml" manifest:media-type="text/xml"/>
</manifest:manifest>"""
    content_xml = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0" office:version="1.2">
  <office:body>
    <office:text>
      <text:p>Liste de courses</text:p>
      <text:p text:bullet-name="•">Pommes</text:p>
      <text:p text:bullet-name="•">Oranges</text:p>
      <text:p text:bullet-name="•">Bananes</text:p>
      <text:p>Fin de la liste</text:p>
    </office:text>
  </office:body>
</office:document-content>"""

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", mimetype)
        zf.writestr("META-INF/manifest.xml", manifest)
        zf.writestr("content.xml", content_xml)
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_odt_with_special_chars_bytes() -> bytes:
    """Generate an ODT document with special characters using raw XML (zip-based)."""
    import zipfile

    mimetype = "application/vnd.oasis.opendocument.text"
    manifest = """<?xml version="1.0" encoding="UTF-8"?>
<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0" manifest:version="1.2">
  <manifest:file-entry manifest:full-path="/" manifest:media-type="application/vnd.oasis.opendocument.text"/>
  <manifest:file-entry manifest:full-path="content.xml" manifest:media-type="text/xml"/>
</manifest:manifest>"""
    content_xml = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" office:version="1.2">
  <office:body>
    <office:text>
      <text:p>Caractères spéciaux</text:p>
      <text:p>Àéîôù Ñ ü ö ä</text:p>
      <text:p>Symboles: © ® ™ € £ ¥</text:p>
    </office:text>
  </office:body>
</office:document-content>"""

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", mimetype)
        zf.writestr("META-INF/manifest.xml", manifest)
        zf.writestr("content.xml", content_xml)
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_xlsx_bytes() -> bytes:
    """Generate a minimal XLSX in memory with 1 sheet and a few cells."""
    from openpyxl import Workbook

    wb: Workbook = Workbook()
    ws = wb.active  # type: ignore[misc]
    ws.title = "Feuille1"  # type: ignore[attr-defined]
    ws["A1"] = "Nom"  # type: ignore[index]
    ws["B1"] = "Valeur"  # type: ignore[index]
    ws["A2"] = "A"  # type: ignore[index]
    ws["B2"] = 1  # type: ignore[index]
    ws["A3"] = "B"  # type: ignore[index]
    ws["B3"] = 2  # type: ignore[index]

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_xlsx_multi_sheet_bytes() -> bytes:
    """Generate an XLSX with multiple sheets."""
    from openpyxl import Workbook

    wb: Workbook = Workbook()
    ws1 = wb.active  # type: ignore[misc]
    ws1.title = "Données"  # type: ignore[attr-defined]
    ws1["A1"] = "Produit"  # type: ignore[index]
    ws1["B1"] = "Prix"  # type: ignore[index]
    ws1["A2"] = "Pomme"  # type: ignore[index]
    ws1["B2"] = 1.5  # type: ignore[index]
    ws1["A3"] = "Orange"  # type: ignore[index]
    ws1["B3"] = 2.0  # type: ignore[index]

    ws2 = wb.create_sheet(title="Résumé")  # type: ignore[misc]
    ws2["A1"] = "Total"  # type: ignore[index]
    ws2["B1"] = 3.5  # type: ignore[index]

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_xlsx_merged_cells_bytes() -> bytes:
    """Generate an XLSX with merged cells."""
    from openpyxl import Workbook

    wb: Workbook = Workbook()
    ws = wb.active  # type: ignore[misc]
    ws.title = "Fusionné"  # type: ignore[attr-defined]
    ws["A1"] = "En-tête fusionné"  # type: ignore[index]
    ws.merge_cells("A1:C1")  # type: ignore[attr-defined]
    ws["A2"] = "Col1"  # type: ignore[index]
    ws["B2"] = "Col2"  # type: ignore[index]
    ws["C2"] = "Col3"  # type: ignore[index]
    ws["A3"] = "v1"  # type: ignore[index]
    ws["B3"] = "v2"  # type: ignore[index]
    ws["C3"] = "v3"  # type: ignore[index]

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_xlsx_dates_numbers_bytes() -> bytes:
    """Generate an XLSX with dates and numbers."""
    from datetime import date, datetime

    from openpyxl import Workbook

    wb: Workbook = Workbook()
    ws = wb.active  # type: ignore[misc]
    ws.title = "Types"  # type: ignore[attr-defined]
    ws["A1"] = "Texte"  # type: ignore[index]
    ws["B1"] = "Nombre"  # type: ignore[index]
    ws["C1"] = "Date"  # type: ignore[index]
    ws["D1"] = "Booléen"  # type: ignore[index]
    ws["A2"] = "hello"  # type: ignore[index]
    ws["B2"] = 42.5  # type: ignore[index]
    ws["C2"] = date(2024, 1, 15)  # type: ignore[index]
    ws["D2"] = True  # type: ignore[index]
    ws["A3"] = "world"  # type: ignore[index]
    ws["B3"] = -10  # type: ignore[index]
    ws["C3"] = datetime(2024, 6, 30, 12, 0, 0)  # type: ignore[index]
    ws["D3"] = False  # type: ignore[index]

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_xlsx_special_chars_bytes() -> bytes:
    """Generate an XLSX with Markdown special characters."""
    from openpyxl import Workbook

    wb: Workbook = Workbook()
    ws = wb.active  # type: ignore[misc]
    ws.title = "Spécial"  # type: ignore[attr-defined]
    ws["A1"] = "Colonne A"  # type: ignore[index]
    ws["B1"] = "Colonne B"  # type: ignore[index]
    ws["A2"] = "Texte avec | pipe"  # type: ignore[index]
    ws["B2"] = "Texte avec \\ backslash"  # type: ignore[index]
    ws["A3"] = "Ligne 1\nLigne 2"  # type: ignore[index]

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_xlsx_empty_sheet_bytes() -> bytes:
    """Generate an XLSX with an empty sheet."""
    from openpyxl import Workbook

    wb: Workbook = Workbook()
    ws = wb.active  # type: ignore[misc]
    ws.title = "Vide"  # type: ignore[attr-defined]

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_bytes() -> bytes:
    """Generate a minimal PPTX in memory using python-pptx."""
    from pptx import Presentation  # type: ignore[misc]

    prs: Presentation = Presentation()  # type: ignore[misc]

    # Slide 1: title slide
    slide_layout = prs.slide_layouts[0]  # type: ignore[attr-defined]  # type: ignore[attr-defined]
    slide = prs.slides.add_slide(slide_layout)  # type: ignore[attr-defined]
    slide.shapes.title.text = "Présentation de Test"  # type: ignore[attr-defined]
    slide.placeholders[1].text = "Sous-titre de la présentation"  # type: ignore[attr-defined]

    # Slide 2: title + content with bullets
    slide_layout = prs.slide_layouts[1]  # type: ignore[attr-defined]
    slide = prs.slides.add_slide(slide_layout)  # type: ignore[attr-defined]
    slide.shapes.title.text = "Deuxième Slide"  # type: ignore[attr-defined]
    body_shape = slide.placeholders[1]  # type: ignore[attr-defined]
    tf = body_shape.text_frame  # type: ignore[attr-defined]
    tf.text = "Premier point"  # type: ignore[attr-defined]
    p = tf.add_paragraph()  # type: ignore[attr-defined]
    p.text = "Deuxième point"  # type: ignore[attr-defined]
    p.level = 0
    p = tf.add_paragraph()  # type: ignore[attr-defined]
    p.text = "Sous-point"  # type: ignore[attr-defined]
    p.level = 1

    buf = io.BytesIO()
    prs.save(buf)  # type: ignore[attr-defined]
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_with_table_bytes() -> bytes:
    """Generate a PPTX with a table in memory using python-pptx."""
    from pptx import Presentation  # type: ignore[misc]
    from pptx.util import Inches  # type: ignore[misc]

    prs = Presentation()  # type: ignore[misc]

    # Slide with table
    slide_layout = prs.slide_layouts[5]  # type: ignore[attr-defined]
    slide = prs.slides.add_slide(slide_layout)  # type: ignore[attr-defined]

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))  # type: ignore[attr-defined]
    title_box.text = "Slide avec Tableau"  # type: ignore[attr-defined]

    # Add table
    rows = 3
    cols = 3
    left = Inches(1)
    top = Inches(1)
    width = Inches(6)
    height = Inches(2)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)  # type: ignore[attr-defined]
    table = table_shape.table  # type: ignore[attr-defined]

    table.cell(0, 0).text = "En-tête 1"  # type: ignore[attr-defined]
    table.cell(0, 1).text = "En-tête 2"  # type: ignore[attr-defined]
    table.cell(0, 2).text = "En-tête 3"  # type: ignore[attr-defined]
    table.cell(1, 0).text = "A"  # type: ignore[attr-defined]
    table.cell(1, 1).text = "B"  # type: ignore[attr-defined]
    table.cell(1, 2).text = "C"  # type: ignore[attr-defined]
    table.cell(2, 0).text = "1"  # type: ignore[attr-defined]
    table.cell(2, 1).text = "2"  # type: ignore[attr-defined]
    table.cell(2, 2).text = "3"  # type: ignore[attr-defined]

    buf = io.BytesIO()
    prs.save(buf)  # type: ignore[attr-defined]
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_empty_slide_bytes() -> bytes:
    """Generate a PPTX with an empty slide in memory using python-pptx."""
    from pptx import Presentation  # type: ignore[misc]

    prs = Presentation()  # type: ignore[misc]
    slide_layout = prs.slide_layouts[6]  # type: ignore[attr-defined]
    prs.slides.add_slide(slide_layout)  # type: ignore[attr-defined]
    # No content added — empty slide

    buf = io.BytesIO()
    prs.save(buf)  # type: ignore[attr-defined]
    buf.seek(0)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# ODS Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_ods_bytes() -> bytes:
    """Generate a minimal ODS file with one sheet and some data using odfpy."""
    from odf import table, text  # type: ignore[attr-defined]
    from odf.opendocument import OpenDocumentSpreadsheet  # type: ignore[attr-defined]

    doc = OpenDocumentSpreadsheet()  # type: ignore[attr-defined]
    table_elem = table.Table(name="Feuille1")  # type: ignore[attr-defined]

    # Row 1: headers
    row1 = table.TableRow()  # type: ignore[attr-defined]
    c1 = table.TableCell()  # type: ignore[attr-defined]
    p1 = text.P()  # type: ignore[attr-defined]
    p1.addText("Nom")  # type: ignore[attr-defined]
    c1.addElement(p1)  # type: ignore[attr-defined]
    row1.addElement(c1)  # type: ignore[attr-defined]
    c2 = table.TableCell()  # type: ignore[attr-defined]
    p2 = text.P()  # type: ignore[attr-defined]
    p2.addText("Valeur")  # type: ignore[attr-defined]
    c2.addElement(p2)  # type: ignore[attr-defined]
    table_elem.addElement(row1)  # type: ignore[attr-defined]

    # Row 2: data
    row2 = table.TableRow()  # type: ignore[attr-defined]
    c3 = table.TableCell()  # type: ignore[attr-defined]
    p3 = text.P()  # type: ignore[attr-defined]
    p3.addText("A")  # type: ignore[attr-defined]
    c3.addElement(p3)  # type: ignore[attr-defined]
    row2.addElement(c3)  # type: ignore[attr-defined]
    c4 = table.TableCell()  # type: ignore[attr-defined]
    p4 = text.P()  # type: ignore[attr-defined]
    p4.addText("1")  # type: ignore[attr-defined]
    c4.addElement(p4)  # type: ignore[attr-defined]
    row2.addElement(c4)  # type: ignore[attr-defined]
    table_elem.addElement(row2)  # type: ignore[attr-defined]

    # Row 3: data
    row3 = table.TableRow()  # type: ignore[attr-defined]
    c5 = table.TableCell()  # type: ignore[attr-defined]
    p5 = text.P()  # type: ignore[attr-defined]
    p5.addText("B")  # type: ignore[attr-defined]
    c5.addElement(p5)  # type: ignore[attr-defined]
    row3.addElement(c5)  # type: ignore[attr-defined]
    c6 = table.TableCell()  # type: ignore[attr-defined]
    p6 = text.P()  # type: ignore[attr-defined]
    p6.addText("2")  # type: ignore[attr-defined]
    c6.addElement(p6)  # type: ignore[attr-defined]
    row3.addElement(c6)  # type: ignore[attr-defined]
    table_elem.addElement(row3)  # type: ignore[attr-defined]

    doc.spreadsheet.addElement(table_elem)  # type: ignore[attr-defined]

    buf = io.BytesIO()
    doc.save(buf)  # type: ignore[attr-defined]
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_ods_multi_sheet_bytes() -> bytes:
    """Generate an ODS file with multiple sheets using odfpy."""
    from odf import table, text  # type: ignore[attr-defined]
    from odf.opendocument import OpenDocumentSpreadsheet  # type: ignore[attr-defined]

    doc = OpenDocumentSpreadsheet()  # type: ignore[attr-defined]

    # Sheet 1: Données
    sheet1 = table.Table(name="Données")  # type: ignore[attr-defined]
    row1 = table.TableRow()  # type: ignore[attr-defined]
    c1 = table.TableCell()  # type: ignore[attr-defined]
    p1 = text.P()  # type: ignore[attr-defined]
    p1.addText("Produit")  # type: ignore[attr-defined]
    c1.addElement(p1)  # type: ignore[attr-defined]
    row1.addElement(c1)  # type: ignore[attr-defined]
    c2 = table.TableCell()  # type: ignore[attr-defined]
    p2 = text.P()  # type: ignore[attr-defined]
    p2.addText("Prix")  # type: ignore[attr-defined]
    c2.addElement(p2)  # type: ignore[attr-defined]
    sheet1.addElement(row1)  # type: ignore[attr-defined]

    row2 = table.TableRow()  # type: ignore[attr-defined]
    c3 = table.TableCell()  # type: ignore[attr-defined]
    p3 = text.P()  # type: ignore[attr-defined]
    p3.addText("Pomme")  # type: ignore[attr-defined]
    c3.addElement(p3)  # type: ignore[attr-defined]
    row2.addElement(c3)  # type: ignore[attr-defined]
    c4 = table.TableCell()  # type: ignore[attr-defined]
    p4 = text.P()  # type: ignore[attr-defined]
    p4.addText("1.5")  # type: ignore[attr-defined]
    c4.addElement(p4)  # type: ignore[attr-defined]
    sheet1.addElement(row2)  # type: ignore[attr-defined]

    row3 = table.TableRow()  # type: ignore[attr-defined]
    c5 = table.TableCell()  # type: ignore[attr-defined]
    p5 = text.P()  # type: ignore[attr-defined]
    p5.addText("Orange")  # type: ignore[attr-defined]
    c5.addElement(p5)  # type: ignore[attr-defined]
    row3.addElement(c5)  # type: ignore[attr-defined]
    c6 = table.TableCell()  # type: ignore[attr-defined]
    p6 = text.P()  # type: ignore[attr-defined]
    p6.addText("2.0")  # type: ignore[attr-defined]
    c6.addElement(p6)  # type: ignore[attr-defined]
    sheet1.addElement(row3)  # type: ignore[attr-defined]

    doc.spreadsheet.addElement(sheet1)  # type: ignore[attr-defined]

    # Sheet 2: Résumé
    sheet2 = table.Table(name="Résumé")  # type: ignore[attr-defined]
    row4 = table.TableRow()  # type: ignore[attr-defined]
    c7 = table.TableCell()  # type: ignore[attr-defined]
    p7 = text.P()  # type: ignore[attr-defined]
    p7.addText("Total")  # type: ignore[attr-defined]
    c7.addElement(p7)  # type: ignore[attr-defined]
    c8 = table.TableCell()  # type: ignore[attr-defined]
    p8 = text.P()  # type: ignore[attr-defined]
    p8.addText("3.5")  # type: ignore[attr-defined]
    c8.addElement(p8)  # type: ignore[attr-defined]
    sheet2.addElement(row4)  # type: ignore[attr-defined]

    doc.spreadsheet.addElement(sheet2)  # type: ignore[attr-defined]

    buf = io.BytesIO()
    doc.save(buf)  # type: ignore[attr-defined]
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_ods_empty_sheet_bytes() -> bytes:
    """Generate an ODS file with an empty sheet using odfpy."""
    # type: ignore[attr-defined]
    from odf import table  # type: ignore[attr-defined]
    from odf.opendocument import OpenDocumentSpreadsheet  # type: ignore[attr-defined]

    doc = OpenDocumentSpreadsheet()  # type: ignore[attr-defined]
    table_elem = table.Table(name="Vide")  # type: ignore[attr-defined]
    # No rows added — empty sheet
    doc.spreadsheet.addElement(table_elem)  # type: ignore[attr-defined]

    buf = io.BytesIO()
    doc.save(buf)  # type: ignore[attr-defined]
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_ods_special_chars_bytes() -> bytes:
    """Generate an ODS file with special characters using odfpy."""
    # type: ignore[attr-defined]
    from odf import table, text  # type: ignore[attr-defined]
    from odf.opendocument import OpenDocumentSpreadsheet  # type: ignore[attr-defined]

    doc = OpenDocumentSpreadsheet()  # type: ignore[attr-defined]
    table_elem = table.Table(name="Spécial")  # type: ignore[attr-defined]

    # Row 1: headers
    row1 = table.TableRow()  # type: ignore[attr-defined]
    c1 = table.TableCell()  # type: ignore[attr-defined]
    p1 = text.P()  # type: ignore[attr-defined]
    p1.addText("Colonne A")  # type: ignore[attr-defined]
    c1.addElement(p1)  # type: ignore[attr-defined]
    row1.addElement(c1)  # type: ignore[attr-defined]
    c2 = table.TableCell()  # type: ignore[attr-defined]
    p2 = text.P()  # type: ignore[attr-defined]
    p2.addText("Colonne B")  # type: ignore[attr-defined]
    c2.addElement(p2)  # type: ignore[attr-defined]
    row1.addElement(c2)  # type: ignore[attr-defined]
    table_elem.addElement(row1)  # type: ignore[attr-defined]

    # Row 2: pipe and backslash
    row2 = table.TableRow()  # type: ignore[attr-defined]
    c3 = table.TableCell()  # type: ignore[attr-defined]
    p3 = text.P()  # type: ignore[attr-defined]
    p3.addText("Texte avec | pipe")  # type: ignore[attr-defined]
    c3.addElement(p3)  # type: ignore[attr-defined]
    row2.addElement(c3)  # type: ignore[attr-defined]
    c4 = table.TableCell()  # type: ignore[attr-defined]
    p4 = text.P()  # type: ignore[attr-defined]
    p4.addText("Texte avec \\ backslash")  # type: ignore[attr-defined]
    c4.addElement(p4)  # type: ignore[attr-defined]
    row2.addElement(c4)  # type: ignore[attr-defined]
    table_elem.addElement(row2)  # type: ignore[attr-defined]

    # Row 3: accents
    row3 = table.TableRow()  # type: ignore[attr-defined]
    c5 = table.TableCell()  # type: ignore[attr-defined]
    p5 = text.P()  # type: ignore[attr-defined]
    p5.addText("Àéîôù")  # type: ignore[attr-defined]
    c5.addElement(p5)  # type: ignore[attr-defined]
    row3.addElement(c5)  # type: ignore[attr-defined]
    c6 = table.TableCell()  # type: ignore[attr-defined]
    p6 = text.P()  # type: ignore[attr-defined]
    p6.addText("Ñ ü ö ä")  # type: ignore[attr-defined]
    c6.addElement(p6)  # type: ignore[attr-defined]
    row3.addElement(c6)  # type: ignore[attr-defined]
    table_elem.addElement(row3)  # type: ignore[attr-defined]

    doc.spreadsheet.addElement(table_elem)  # type: ignore[attr-defined]

    buf = io.BytesIO()
    doc.save(buf)  # type: ignore[attr-defined]
    buf.seek(0)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# ODP Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_odp_bytes() -> bytes:
    """Generate a minimal ODP file with slides and text using odfpy."""
    # type: ignore[misc]
    from odf import draw  # type: ignore[misc]
    from odf import opendocument  # type: ignore[misc]
    from odf import text  # type: ignore[misc]

    doc = opendocument.OpenDocumentPresentation()  # type: ignore[misc]

    # Slide 1: title slide
    page1 = draw.Page(name="Slide 1", masterpagename="Default")  # type: ignore[misc]
    title_frame = draw.Frame(name="TitleFrame")  # type: ignore[misc]
    title_text_box = draw.TextBox()  # type: ignore[misc]
    title_para = text.P()  # type: ignore[misc]
    title_para.addText("Présentation de Test")  # type: ignore[attr-defined]
    title_text_box.addElement(title_para)  # type: ignore[attr-defined]
    title_frame.addElement(title_text_box)  # type: ignore[attr-defined]
    page1.addElement(title_frame)  # type: ignore[attr-defined]

    # Subtitle frame
    subtitle_frame = draw.Frame(name="SubtitleFrame")  # type: ignore[misc]
    subtitle_text_box = draw.TextBox()  # type: ignore[misc]
    subtitle_para = text.P()  # type: ignore[misc]
    subtitle_para.addText("Sous-titre de la présentation")  # type: ignore[attr-defined]
    subtitle_text_box.addElement(subtitle_para)  # type: ignore[attr-defined]
    subtitle_frame.addElement(subtitle_text_box)  # type: ignore[attr-defined]
    page1.addElement(subtitle_frame)  # type: ignore[attr-defined]

    doc.presentation.addElement(page1)  # type: ignore[attr-defined]

    # Slide 2: content slide
    page2 = draw.Page(name="Slide 2", masterpagename="Default")  # type: ignore[misc]
    title_frame2 = draw.Frame(name="TitleFrame2")  # type: ignore[misc]
    title_text_box2 = draw.TextBox()  # type: ignore[misc]
    title_para2 = text.P()  # type: ignore[misc]
    title_para2.addText("Deuxième Slide")  # type: ignore[attr-defined]
    title_text_box2.addElement(title_para2)  # type: ignore[attr-defined]
    title_frame2.addElement(title_text_box2)  # type: ignore[attr-defined]
    page2.addElement(title_frame2)  # type: ignore[attr-defined]

    # Content frame with list items
    content_frame = draw.Frame(name="ContentFrame")  # type: ignore[misc]
    
    # Content text box with bullet points
    content_text_box = draw.TextBox()  # type: ignore[misc]
    para1 = text.P()  # type: ignore[misc]
    para1.addText("Premier point")  # type: ignore[attr-defined]
    content_text_box.addElement(para1)  # type: ignore[attr-defined]
    para2 = text.P()  # type: ignore[misc]
    para2.addText("Deuxième point")  # type: ignore[attr-defined]
    content_text_box.addElement(para2)  # type: ignore[attr-defined]
    para3 = text.P()  # type: ignore[misc]
    para3.addText("Troisième point")  # type: ignore[attr-defined]
    content_text_box.addElement(para3)  # type: ignore[attr-defined]
    content_frame.addElement(content_text_box)  # type: ignore[attr-defined]
    page2.addElement(content_frame)  # type: ignore[attr-defined]

    doc.presentation.addElement(page2)  # type: ignore[attr-defined]

    buf = io.BytesIO()
    doc.save(buf)  # type: ignore[attr-defined]
    buf.seek(0)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# JSON/JSONL Test Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def expected_json_output_pdf() -> Dict[str, Any]:
    """Expected JSON output for PDF conversion (Hello World document)."""
    return {
        "format": "pdf",
        "pages": [
            {
                "page_idx": 0,
                "markdown_text": "# Hello World\n\nLigne deux",
                "links": []
            }
        ],
        "metadata": {"source_format": "pdf", "num_pages": 1},
        "success": True,
        "timestamp": None  # Will be set by actual conversion
    }


@pytest.fixture
def expected_json_output_docx() -> Dict[str, Any]:
    """Expected JSON output for DOCX conversion."""
    return {
        "format": "docx",
        "pages": [
            {
                "page_idx": 0,
                "markdown_text": "# Bonjour le monde\n\n## Deuxième paragraphe",
                "links": []
            },
            {
                "page_idx": 1,
                "markdown_text": "# Deuxième paragraphe",
                "links": []
            }
        ],
        "metadata": {"source_format": "docx", "num_pages": 2},
        "success": True,
        "timestamp": None
    }


@pytest.fixture
def expected_json_output_odt() -> Dict[str, Any]:
    """Expected JSON output for ODT conversion."""
    return {
        "format": "odt",
        "pages": [
            {
                "page_idx": 0,
                "markdown_text": "# Titre du document\n\n## Premier paragraphe\n\n## Deuxième paragraphe",
                "links": []
            }
        ],
        "metadata": {"source_format": "odt", "num_pages": 1},
        "success": True,
        "timestamp": None
    }


@pytest.fixture
def expected_json_output_xlsx() -> Dict[str, Any]:
    """Expected JSON output for XLSX conversion."""
    return {
        "format": "xlsx",
        "sheets": [
            {
                "name": "Feuille1",
                "data": [
                    ["Nom", "Valeur"],
                    ["A", 1],
                    ["B", 2]
                ],
                "headers": ["Nom", "Valeur"]
            }
        ],
        "metadata": {"source_format": "xlsx", "num_sheets": 1},
        "success": True,
        "timestamp": None
    }


@pytest.fixture
def expected_json_output_pptx() -> Dict[str, Any]:
    """Expected JSON output for PPTX conversion."""
    return {
        "format": "pptx",
        "slides": [
            {
                "index": 0,
                "title": "Présentation de Test",
                "content": ["Sous-titre de la présentation"],
                "tables": []
            },
            {
                "index": 1,
                "title": "Deuxième Slide",
                "content": ["Premier point", "Deuxième point", "Sous-point"],
                "tables": []
            }
        ],
        "metadata": {"source_format": "pptx", "num_slides": 2},
        "success": True,
        "timestamp": None
    }


@pytest.fixture
def expected_json_output_html() -> Dict[str, Any]:
    """Expected JSON output for HTML conversion."""
    return {
        "format": "html",
        "pages": [
            {
                "page_idx": 0,
                "markdown_text": "# Main Title\n\nThis is a **bold** and *italic* paragraph.\n\nVisit our website for more info.",
                "links": ["https://example.com"]
            }
        ],
        "metadata": {"source_format": "html", "num_pages": 1},
        "success": True,
        "timestamp": None
    }


@pytest.fixture
def expected_json_output_ods() -> Dict[str, Any]:
    """Expected JSON output for ODS conversion."""
    return {
        "format": "ods",
        "sheets": [
            {
                "name": "Feuille1",
                "data": [
                    ["Nom", "Valeur"],
                    ["A", 1],
                    ["B", 2]
                ],
                "headers": ["Nom", "Valeur"]
            }
        ],
        "metadata": {"source_format": "ods", "num_sheets": 1},
        "success": True,
        "timestamp": None
    }


@pytest.fixture
def expected_json_output_odp() -> Dict[str, Any]:
    """Expected JSON output for ODP conversion."""
    return {
        "format": "odp",
        "slides": [
            {
                "index": 0,
                "title": "Présentation de Test",
                "content": ["Sous-titre de la présentation"],
                "lists": []
            },
            {
                "index": 1,
                "title": "Deuxième Slide",
                "content": ["Premier point", "Deuxième point", "Troisième point"],
                "lists": []
            }
        ],
        "metadata": {"source_format": "odp", "num_slides": 2},
        "success": True,
        "timestamp": None
    }


@pytest.fixture
def expected_jsonl_output_pdf() -> List[str]:
    """Expected JSONL output for PDF conversion (one event per page)."""
    return [
        '{"type": "start", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}',
        '{"type": "chunk", "page_idx": 0, "markdown_text": "# Hello World\\n\\nLigne deux", "links": [], "offset": 0, "length": 28}',
        '{"type": "end", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}'
    ]


@pytest.fixture
def expected_jsonl_output_docx() -> List[str]:
    """Expected JSONL output for DOCX conversion."""
    return [
        '{"type": "start", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}',
        '{"type": "chunk", "page_idx": 0, "markdown_text": "# Bonjour le monde\\n\\n## Deuxième paragraphe", "links": [], "offset": 0, "length": 42}',
        '{"type": "chunk", "page_idx": 1, "markdown_text": "# Deuxième paragraphe", "links": [], "offset": 42, "length": 23}',
        '{"type": "end", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}'
    ]


@pytest.fixture
def expected_jsonl_output_odt() -> List[str]:
    """Expected JSONL output for ODT conversion."""
    return [
        '{"type": "start", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}',
        '{"type": "chunk", "page_idx": 0, "markdown_text": "# Titre du document\\n\\n## Premier paragraphe\\n\\n## Deuxième paragraphe", "links": [], "offset": 0, "length": 72}',
        '{"type": "end", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}'
    ]


@pytest.fixture
def expected_jsonl_output_xlsx() -> List[str]:
    """Expected JSONL output for XLSX conversion."""
    return [
        '{"type": "start", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}',
        '{"type": "chunk", "page_idx": 0, "markdown_text": [["Nom","Valeur"],["A",1.0],"[B",2.0]]", "links": [], "offset": 0, "length": 35}',
        '{"type": "end", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}'
    ]


@pytest.fixture
def expected_jsonl_output_pptx() -> List[str]:
    """Expected JSONL output for PPTX conversion."""
    return [
        '{"type": "start", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}',
        '{"type": "chunk", "page_idx": 0, "markdown_text": "Présentation de Test\\nSous-titre de la présentation", "links": [], "offset": 0, "length": 48}',
        '{"type": "chunk", "page_idx": 1, "markdown_text": "Deuxième Slide\\nPremier point\\nDeuxième point\\nSous-point", "links": [], "offset": 48, "length": 59}',
        '{"type": "end", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}'
    ]


@pytest.fixture
def expected_jsonl_output_html() -> list[str]:
    """Expected JSONL output for HTML conversion."""
    return [
        '{"type": "start", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}',
        '{"type": "chunk", "page_idx": 0, "markdown_text": "# Main Title\\nThis is a **bold** and *italic* paragraph.\\nVisit our website for more info.", "links": ["https://example.com"], "offset": 0, "length": 95}',
        '{"type": "end", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}'
    ]


@pytest.fixture
def expected_jsonl_output_ods() -> list[str]:
    """Expected JSONL output for ODS conversion."""
    return [
        '{"type": "start", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}',
        '{"type": "chunk", "page_idx": 0, "markdown_text": [["Nom","Valeur"],["A",1.0],"[B",2.0]]", "links": [], "offset": 0, "length": 35}',
        '{"type": "end", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}'
    ]


@pytest.fixture
def expected_jsonl_output_odp() -> list[str]:
    """Expected JSONL output for ODP conversion."""
    return [
        '{"type": "start", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}',
        '{"type": "chunk", "page_idx": 0, "markdown_text": "Présentation de Test\\nSous-titre de la présentation", "links": [], "offset": 0, "length": 48}',
        '{"type": "chunk", "page_idx": 1, "markdown_text": "Deuxième Slide\\nPremier point\\nDeuxième point\\nTroisième point", "links": [], "offset": 48, "length": 62}',
        '{"type": "end", "page_idx": null, "markdown_text": "", "links": [], "offset": 0, "length": 0}'
    ]


# ---------------------------------------------------------------------------
# End of Test Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_odp_with_list_bytes() -> bytes:
    """Generate an ODP file with bullet list items using odfpy."""
    from odf import draw, opendocument, text  # type: ignore[attr-defined]

    doc = opendocument.OpenDocumentPresentation()  # type: ignore[attr-defined]

    page = draw.Page(name="Slide 1", masterpagename="Default")  # type: ignore[attr-defined]
    title_frame = draw.Frame(name="TitleFrame")  # type: ignore[attr-defined]
    title_text_box = draw.TextBox()  # type: ignore[attr-defined]
    title_para = text.P()  # type: ignore[attr-defined]
    title_para.addText("Liste de courses")  # type: ignore[attr-defined]
    title_text_box.addElement(title_para)  # type: ignore[attr-defined]
    title_frame.addElement(title_text_box)  # type: ignore[attr-defined]
    page.addElement(title_frame)  # type: ignore[attr-defined]

    content_frame = draw.Frame(name="ContentFrame")  # type: ignore[attr-defined]
    content_text_box = draw.TextBox()  # type: ignore[attr-defined]
    para1 = text.P()  # type: ignore[attr-defined]
    para1.addText("Pommes")  # type: ignore[attr-defined]
    content_text_box.addElement(para1)  # type: ignore[attr-defined]
    para2 = text.P()  # type: ignore[attr-defined]
    para2.addText("Oranges")  # type: ignore[attr-defined]
    content_text_box.addElement(para2)  # type: ignore[attr-defined]
    para3 = text.P()  # type: ignore[attr-defined]
    para3.addText("Bananes")  # type: ignore[attr-defined]
    content_text_box.addElement(para3)  # type: ignore[attr-defined]
    content_frame.addElement(content_text_box)  # type: ignore[attr-defined]
    page.addElement(content_frame)  # type: ignore[attr-defined]

    doc.presentation.addElement(page)  # type: ignore[attr-defined]

    buf = io.BytesIO()
    doc.save(buf)  # type: ignore[attr-defined]
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_odp_with_special_chars_bytes() -> bytes:
    """Generate an ODP file with special characters using odfpy."""
    from odf import draw, opendocument, text  # type: ignore[attr-defined]

    doc = opendocument.OpenDocumentPresentation()  # type: ignore[attr-defined]

    page = draw.Page(name="Slide 1", masterpagename="Default")  # type: ignore[attr-defined]
    title_frame = draw.Frame(name="TitleFrame")  # type: ignore[attr-defined]
    title_text_box = draw.TextBox()  # type: ignore[attr-defined]
    title_para = text.P()  # type: ignore[attr-defined]
    title_para.addText("Caractères spéciaux")  # type: ignore[attr-defined]
    title_text_box.addElement(title_para)  # type: ignore[attr-defined]
    title_frame.addElement(title_text_box)  # type: ignore[attr-defined]
    page.addElement(title_frame)  # type: ignore[attr-defined]

    content_frame = draw.Frame(name="ContentFrame")  # type: ignore[attr-defined]
    content_text_box = draw.TextBox()  # type: ignore[attr-defined]
    para1 = text.P()  # type: ignore[attr-defined]
    para1.addText("Àéîôù Ñ ü ö ä")  # type: ignore[attr-defined]
    content_text_box.addElement(para1)  # type: ignore[attr-defined]
    para2 = text.P()  # type: ignore[attr-defined]
    para2.addText("Symboles: © ® ™ € £ ¥")  # type: ignore[attr-defined]
    content_text_box.addElement(para2)  # type: ignore[attr-defined]
    content_frame.addElement(content_text_box)  # type: ignore[attr-defined]
    page.addElement(content_frame)  # type: ignore[attr-defined]

    doc.presentation.addElement(page)  # type: ignore[attr-defined]

    buf = io.BytesIO()
    doc.save(buf)  # type: ignore[attr-defined]
    buf.seek(0)
    return buf.getvalue()


@pytest.fixture
def sample_odp_with_groups_bytes() -> bytes:
    """Generate an ODP file with draw:g groups (LibreOffice-style structure).

    LibreOffice wraps frames in draw:g (group) elements, which was previously
    not handled by the ODP converter.
    """
    import zipfile

    content_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
        ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0"'
        ' xmlns:draw="urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"'
        ' xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"'
        ' office:version="1.3">'
        '  <office:body>'
        '    <office:presentation>'
        '      <draw:page draw:name="Page1" draw:master-page-name="Default">'
        '        <draw:g draw:style-name="dp1">'
        '          <draw:frame draw:name="pm-title" draw:style-name="Title">'
        '            <draw:text-frame>'
        '              <text:p>Titre dans un groupe</text:p>'
        '            </draw:text-frame>'
        '          </draw:frame>'
        '        </draw:g>'
        '        <draw:g draw:style-name="dp1">'
        '          <draw:frame draw:name="pm-subtitle">'
        '            <draw:text-frame>'
        '              <text:p>Sous-titre dans un groupe</text:p>'
        '            </draw:text-frame>'
        '          </draw:frame>'
        '        </draw:g>'
        '        <draw:frame draw:name="ContentFrame">'
        '          <draw:text-frame>'
        '            <text:p>Contenu direct</text:p>'
        '          </draw:text-frame>'
        '        </draw:frame>'
        '      </draw:page>'
        '    </office:presentation>'
        '  </office:body>'
        '</office:document-content>'
    )

    manifest_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0">'
        '  <manifest:file-entry manifest:full-path="/"'
        '    manifest:media-type="application/vnd.oasis.opendocument.presentation"/>'
        '  <manifest:file-entry manifest:full-path="content.xml"'
        '    manifest:media-type="text/xml"/>'
        "</manifest:manifest>"
    )

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("content.xml", content_xml)
        zf.writestr("META-INF/manifest.xml", manifest_xml)
    buf.seek(0)
    return buf.getvalue()
